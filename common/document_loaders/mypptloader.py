from langchain.document_loaders.unstructured import UnstructuredFileLoader
from typing import List
import tqdm


class RapidOCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        def ppt2text(filepath):
            from pptx import Presentation
            from PIL import Image
            import numpy as np
            from io import BytesIO
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm.tqdm(total=len(prs.slides),
                               desc="RapidOCRPPTLoader slide index: 1")
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number))
                b_unit.refresh()
                sorted_shapes = sorted(slide.shapes,
                                       key=lambda x: (x.top, x.left))  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        from unstructured.partition.text import partition_text
        return partition_text(text=text, **self.unstructured_kwargs)


if __name__ == '__main__':
    loader = RapidOCRPPTLoader(file_path="../tests/samples/ocr_test.pptx")
    docs = loader.load()
    print(docs)
